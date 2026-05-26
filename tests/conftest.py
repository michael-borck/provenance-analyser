"""Fixture builders — every test artefact is constructed in tmp_path via the
authoring libraries themselves, so there are no binary fixtures on disk."""
from __future__ import annotations

import zipfile
from datetime import datetime, timezone
from pathlib import Path

import pytest


@pytest.fixture
def docx_with_metadata(tmp_path: Path) -> Path:
    """A .docx with author + revision + TotalTime + Words populated."""
    from docx import Document

    doc = Document()
    doc.core_properties.author = "Jane Student"
    doc.core_properties.last_modified_by = "Jane Student"
    doc.core_properties.revision = 4
    doc.add_paragraph("This is the body text. " * 20)
    path = tmp_path / "essay.docx"
    doc.save(path)

    # python-docx doesn't write TotalTime/Words to app.xml, so patch it in via zip
    # (this is what real Word does — we're emulating realistic output for the test).
    _patch_app_xml(path, {"TotalTime": "45", "Words": "100", "Application": "Microsoft Office Word"})
    return path


@pytest.fixture
def docx_low_edit_time(tmp_path: Path) -> Path:
    """Suspicious: 2 minutes of editing for 1000+ words → edit_time_low_for_size flag."""
    from docx import Document

    doc = Document()
    doc.core_properties.author = "Jane"
    doc.add_paragraph("word " * 1500)
    path = tmp_path / "suspect.docx"
    doc.save(path)
    _patch_app_xml(path, {"TotalTime": "2", "Words": "1500", "Application": "Microsoft Office Word"})
    return path


@pytest.fixture
def docx_with_ai_marker(tmp_path: Path) -> Path:
    """A .docx where Application is 'ChatGPT' → ai_generation_marker flag."""
    from docx import Document

    doc = Document()
    doc.add_paragraph("generated text")
    path = tmp_path / "ai.docx"
    doc.save(path)
    _patch_app_xml(path, {"Application": "ChatGPT"})
    return path


@pytest.fixture
def docx_minimal(tmp_path: Path) -> Path:
    """A docx with no author / no metadata → missing_metadata flag."""
    from docx import Document

    doc = Document()
    doc.add_paragraph("body")
    path = tmp_path / "minimal.docx"
    doc.save(path)
    return path


@pytest.fixture
def pptx_with_metadata(tmp_path: Path) -> Path:
    from pptx import Presentation

    prs = Presentation()
    prs.core_properties.author = "Slide Author"
    prs.core_properties.last_modified_by = "Reviewer"  # mismatch → flag
    prs.core_properties.revision = 2
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.slides.add_slide(prs.slide_layouts[1])
    path = tmp_path / "deck.pptx"
    prs.save(path)
    _patch_app_xml(path, {"Application": "Microsoft Office PowerPoint", "TotalTime": "30", "Slides": "2"})
    return path


@pytest.fixture
def xlsx_with_metadata(tmp_path: Path) -> Path:
    from openpyxl import Workbook

    wb = Workbook()
    wb.properties.creator = "Spreadsheet Author"
    wb.properties.lastModifiedBy = "Spreadsheet Author"
    wb.properties.revision = "3"
    wb.active.title = "Sheet1"
    wb.create_sheet("Sheet2")
    path = tmp_path / "model.xlsx"
    wb.save(path)
    _patch_app_xml(path, {"Application": "Microsoft Excel", "TotalTime": "60"})
    return path


@pytest.fixture
def pdf_with_metadata(tmp_path: Path) -> Path:
    """A real .pdf written via reportlab so pypdf can parse it."""
    from reportlab.lib.pagesizes import LETTER
    from reportlab.pdfgen import canvas

    path = tmp_path / "report.pdf"
    c = canvas.Canvas(str(path), pagesize=LETTER)
    c.setAuthor("PDF Author")
    c.setTitle("Test Report")
    c.setSubject("Provenance test")
    c.setCreator("ReportLab")
    c.drawString(72, 720, "Hello PDF")
    c.showPage()
    c.save()
    return path


# ── helpers ──────────────────────────────────────────────────────────────


def _patch_app_xml(path: Path, fields: dict[str, str]) -> None:
    """Rewrite docProps/app.xml inside the OOXML zip to inject extended properties.

    python-docx / python-pptx / openpyxl all leave the extended-properties file
    minimal or absent; tests need realistic TotalTime / Words / Application so
    the extractors have something to read.
    """
    # Read the existing zip into memory, then write a new one.
    EP_NS = "http://schemas.openxmlformats.org/officeDocument/2006/extended-properties"
    VT_NS = "http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes"

    body = ['<?xml version="1.0" encoding="UTF-8" standalone="yes"?>',
            f'<Properties xmlns="{EP_NS}" xmlns:vt="{VT_NS}">']
    for k, v in fields.items():
        body.append(f"<{k}>{v}</{k}>")
    body.append("</Properties>")
    app_xml = "\n".join(body)

    import io
    in_zip_bytes = path.read_bytes()
    out_buf = io.BytesIO()
    with zipfile.ZipFile(io.BytesIO(in_zip_bytes), "r") as src:
        with zipfile.ZipFile(out_buf, "w", zipfile.ZIP_DEFLATED) as dst:
            seen_app = False
            for item in src.infolist():
                if item.filename == "docProps/app.xml":
                    dst.writestr(item, app_xml)
                    seen_app = True
                else:
                    dst.writestr(item, src.read(item.filename))
            if not seen_app:
                dst.writestr("docProps/app.xml", app_xml)
    path.write_bytes(out_buf.getvalue())
